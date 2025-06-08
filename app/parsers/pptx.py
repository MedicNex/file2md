from .base import BaseParser
from loguru import logger
from pptx import Presentation
from app.vision import image_to_markdown
import io
import os

class PptxParser(BaseParser):
    """PPTX文件解析器"""
    
    @classmethod
    def get_supported_extensions(cls) -> list[str]:
        return ['.ppt', '.pptx']
    
    async def parse(self, file_path: str) -> str:
        """解析PPTX文件"""
        try:
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"## 幻灯片 {slide_idx}")
                
                # 提取文本内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # 处理标题和内容
                        if shape.shape_type == 1:  # 标题形状
                            slide_content.append(f"### {text}")
                        else:
                            # 处理列表项
                            lines = text.split('\n')
                            formatted_lines = []
                            for line in lines:
                                line = line.strip()
                                if line:
                                    if line.startswith('•') or line.startswith('-'):
                                        formatted_lines.append(f"- {line[1:].strip()}")
                                    else:
                                        formatted_lines.append(line)
                            slide_content.append('\n'.join(formatted_lines))
                    
                    # 处理图片
                    if hasattr(shape, "image"):
                        try:
                            # 提取图片数据
                            image = shape.image
                            image_bytes = image.blob
                            
                            # 保存为临时文件
                            temp_img_path = self.create_temp_file(suffix='.png', content=image_bytes)
                            
                            # 使用视觉模型解析图片
                            img_markdown = await image_to_markdown(temp_img_path)
                            slide_content.append(img_markdown)
                            
                        except Exception as img_error:
                            logger.warning(f"PPT图片处理失败 幻灯片{slide_idx}: {img_error}")
                            slide_content.append("*图片处理失败*")
                
                if len(slide_content) > 1:  # 有内容才添加
                    content_parts.append('\n\n'.join(slide_content))
            
            markdown_content = '\n\n---\n\n'.join(content_parts)
            
            logger.info(f"成功解析PPTX文件: {file_path}")
            return markdown_content or "PowerPoint文件为空或无法提取内容"
            
        except Exception as e:
            logger.error(f"解析PPTX文件失败 {file_path}: {e}")
            raise Exception(f"PPTX文件解析错误: {str(e)}") 